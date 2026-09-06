"""Patch Lesson 1 v2 PPTX to use Bangla textbook-facing terminology.

Purpose:
- Preserve the v2 visual-engagement/staged-reveal deck.
- Replace instructional English with Bangla textbook language.
- Clarify glass-slab prediction/reveal slides.
- Rename the 'Live demo' slide as a classroom demonstration with a static fallback.

Input expected:
  Class8_Science_Ch11_Lesson1_v2_Visual_Engagement_PPTX_SAFE.pptx

Output:
  Class8_Science_Ch11_Lesson1_v2_Bangla_Textbook_PPTX_SAFE.pptx

Note:
This patch is a reproducibility record of the user-requested language/clarity correction.
The classroom-safe final file should still be LibreOffice-normalized and PowerPoint-smoke-tested.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

SRC = Path("/mnt/data/science_slides_l1_v2/Class8_Science_Ch11_Lesson1_v2_Visual_Engagement_PPTX_SAFE.pptx")
OUT_PRENORM = Path("/mnt/data/science_slides_l1_v2/Class8_Science_Ch11_Lesson1_v2_Bangla_Textbook_PRENORM.pptx")

prs = Presentation(SRC)

# Broad but controlled replacements for student-facing text.
REPLACEMENTS = [
    ("Hook", "হুক"),
    ("Think → Pair → Answer", "ভাবো → পাশে আলোচনা → উত্তর দাও"),
    ("ক্লিক/Next: পরের ধাপ", "ক্লিক/পরের ধাপ"),
    ("light path", "আলোর পথ"),
    ("Step 1", "ধাপ ১"),
    ("Step 2", "ধাপ ২"),
    ("Step 3", "ধাপ ৩"),
    ("source", "উৎস"),
    ("Source", "উৎস"),
    ("lesson", "পাঠ"),
    ("ray model", "রশ্মি-ধারণা"),
    ("pathway", "ধারাবাহিকতা"),
    ("image", "প্রতিবিম্ব"),
    ("focus", "মূল লক্ষ্য"),
    ("Prediction", "অনুমান"),
    ("Transparent / Translucent / Opaque", "স্বচ্ছ / অর্ধস্বচ্ছ / অস্বচ্ছ"),
    ("Transparent", "স্বচ্ছ"),
    ("Translucent", "অর্ধস্বচ্ছ"),
    ("Opaque", "অস্বচ্ছ"),
    ("Live demo", "শ্রেণিকক্ষ প্রদর্শনী"),
    ("Explain", "ব্যাখ্যা"),
    ("observe", "পর্যবেক্ষণ"),
    ("Observe", "পর্যবেক্ষণ"),
    ("Static fallback", "স্থির বিকল্প"),
    ("live demo", "বাস্তব প্রদর্শনী"),
    ("observation", "পর্যবেক্ষণ"),
    ("Ray diagram", "রশ্মিচিত্র"),
    ("boundary", "বিভেদতল"),
    ("interface", "বিভেদতল"),
    ("Interface", "বিভেদতল"),
    ("surface", "পৃষ্ঠ"),
    ("Angles", "কোণ"),
    ("Incident ray", "আপতিত রশ্মি"),
    ("Refracted ray", "প্রতিসরিত রশ্মি"),
    ("Medium 1", "প্রথম মাধ্যম"),
    ("Medium 2", "দ্বিতীয় মাধ্যম"),
    ("medium", "মাধ্যম"),
    ("Medium", "মাধ্যম"),
    ("ray", "রশ্মি"),
    ("Ray", "রশ্মি"),
    ("reflection", "প্রতিফলন"),
    ("Reflection", "প্রতিফলন"),
    ("Refraction", "প্রতিসরণ"),
    ("refraction", "প্রতিসরণ"),
    ("normal", "অভিলম্ব"),
    ("Normal", "অভিলম্ব"),
    ("speed", "বেগ"),
    ("path", "পথ"),
    ("optical density", "আলোকীয় ঘনত্ব"),
    ("Optical density", "আলোকীয় ঘনত্ব"),
    ("mass density", "ভরঘনত্ব"),
    ("direction", "দিক"),
    ("Direction", "দিক"),
    ("Vote", "ভোট"),
    ("YES / NO", "হ্যাঁ / না"),
    ("Exception", "ব্যতিক্রম"),
    ("Rule", "নিয়ম"),
    ("Draw with finger", "আঙুল দিয়ে দেখাও"),
    ("chain", "ধারাবাহিকতা"),
    ("angle", "কোণ"),
    ("Teacher note", "শিক্ষক নির্দেশনা"),
    ("node", "ধাপ"),
    ("slide", "স্লাইড"),
    ("Exit Check", "শেষ যাচাই"),
    ("Exit check", "শেষ যাচাই"),
    ("Teacher click", "শিক্ষক ক্লিক"),
    ("FLEX", "ঐচ্ছিক"),
    ("Glass slab", "কাচের পাত"),
    ("Parallel-sided glass slab", "সমান্তরাল-পৃষ্ঠ কাচের পাত"),
    ("Parallel slab", "সমান্তরাল-পৃষ্ঠ কাচের পাত"),
    ("emergent ray", "নির্গত রশ্মি"),
    ("incident ray", "আপতিত রশ্মি"),
    ("parallel", "সমান্তরাল"),
    ("Parallel", "সমান্তরাল"),
    ("Guardrail", "সতর্কতা"),
    ("equality", "সমতা"),
    ("lateral shift", "পার্শ্ব সরণ"),
    ("Predict before reveal", "প্রকাশের আগে অনুমান"),
    ("Simulation", "সিমুলেশন"),
    ("incidence angle", "আপতন কোণ"),
    ("Primary", "প্রাথমিক"),
    ("offline-friendly support", "অফলাইন-সহায়ক"),
    ("Prediction task", "অনুমান করো"),
    ("Launch", "খুলুন"),
    ("Alternate", "বিকল্প"),
    ("Interactive", "ইন্টারঅ্যাকটিভ"),
    ("CORE fallback", "মূল বিকল্প"),
    ("Internet", "ইন্টারনেট"),
    ("staged diagrams", "ধাপে ধাপে রশ্মিচিত্র"),
]

TARGETED_REPLACEMENTS = {
    "এই পাঠ-এ": "এই পাঠে",
    "রশ্মিচিত্র বানানোর ভাষা": "রশ্মিচিত্র আঁকার ভাষা",
    "ধাপ ১: দুই মাধ্যমের বিভেদতল বা বিভেদতল": "ধাপ ১: দুই মাধ্যমের বিভেদতল",
    "বিভেদতল / বিভেদতল": "বিভেদতল",
    "প্রতিফলন vs প্রতিসরণ": "প্রতিফলন বনাম প্রতিসরণ",
    "রশ্মি অভিলম্ব-এর": "রশ্মি অভিলম্বের",
    "অভিলম্ব-এর": "অভিলম্বের",
    "স্লাইড-এ": "স্লাইডে",
    "বিভেদতল-এ": "বিভেদতলে",
    "মাধ্যম-এ": "মাধ্যমে",
    "water পৃষ্ঠ": "পানির পৃষ্ঠ",
    "bending rules": "প্রতিসরণের নিয়ম",
}

# exact shape overrides for cases split across text runs or created by broad replacements
OVERRIDES = [
    (10, "water পৃষ্ঠ", "পর্যবেক্ষণ: পানির পৃষ্ঠের উপরের ও নিচের অংশ কি একই সরলরেখায় দেখা যায়?"),
    (24, "কোন line", "১. কোণ কোন রেখা থেকে মাপি?"),
    (20, "rarer মাধ্যম", "নিয়ম: বিরল মাধ্যম থেকে ঘন মাধ্যমে গেলে রশ্মি অভিলম্বের দিকে বেঁকে যায়।"),
    (22, "rarer মাধ্যম", "নিয়ম: ঘন মাধ্যম থেকে বিরল মাধ্যমে গেলে রশ্মি অভিলম্ব থেকে দূরে বেঁকে যায়।"),
    (21, "আগের rule", "আগের নিয়ম উল্টে গেলে পথ কী হবে?"),
    (26, "emergent রশ্মি", "প্রকাশের আগে অনুমান: নির্গত রশ্মি কি আপতিত রশ্মির সাথে সমান্তরাল হবে?"),
    (27, "emergent রশ্মি সমান্তরাল", "সমান্তরাল-পৃষ্ঠ কাচের পাত: নির্গত রশ্মি আপতিত রশ্মির সমান্তরাল হতে পারে"),
    (27, "সমান্তরাল-sided slab", "সতর্কতা: এই সমতা সমান্তরাল-পৃষ্ঠ কাচের পাতের জ্যামিতির জন্য, সব বিভেদতলে নয়।"),
    (27, "incident রশ্মি ও emergent", "জ্যামিতিনির্ভর নিয়ম: আপতিত রশ্মি ও নির্গত রশ্মি সমান্তরাল; পার্শ্ব সরণ থাকতে পারে।"),
    (28, "পাঠ complete", "মূল বিকল্প: L1-S09–L1-S12 ধাপে ধাপে রশ্মিচিত্র। ইন্টারনেট না থাকলেও পাঠ সম্পূর্ণ।"),
]


def set_text(shape, text):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Noto Sans Bengali"


def add_text(slide, text, x, y, w, h, size=14, color=(246, 251, 255), bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Noto Sans Bengali"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = RGBColor(*color)
    return box


for slide in prs.slides:
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text
                for src, dst in REPLACEMENTS:
                    text = text.replace(src, dst)
                for src, dst in TARGETED_REPLACEMENTS.items():
                    text = text.replace(src, dst)
                run.text = text

for slide_no, fragment, new_text in OVERRIDES:
    slide = prs.slides[slide_no - 1]
    for shape in slide.shapes:
        if hasattr(shape, "text") and fragment in shape.text:
            set_text(shape, new_text)

# Add clarifying labels to glass-slab slides.
slide = prs.slides[25]
add_text(slide, "কাচের পাত", 5.4, 2.25, 1.6, 0.35, size=16, color=(155, 234, 254), bold=True)
add_text(slide, "আপতিত রশ্মি", 1.55, 2.55, 1.5, 0.3, size=14, color=(255, 209, 102), bold=True)
add_text(slide, "প্রথম বিভেদতল", 3.05, 5.55, 1.8, 0.3, size=12, color=(183, 205, 232))
add_text(slide, "দ্বিতীয় বিভেদতল", 9.1, 5.55, 1.9, 0.3, size=12, color=(183, 205, 232))

slide = prs.slides[26]
add_text(slide, "আপতিত রশ্মি", 1.55, 2.43, 1.5, 0.3, size=13, color=(255, 209, 102), bold=True)
add_text(slide, "প্রতিসরিত রশ্মি", 5.0, 3.42, 1.7, 0.3, size=13, color=(46, 229, 157), bold=True)
add_text(slide, "নির্গত রশ্মি", 10.0, 4.95, 1.6, 0.3, size=13, color=(255, 209, 102), bold=True)
add_text(slide, "পাতলা কমলা রেখা: আপতিত রশ্মির সমান্তরাল সহায়ক রেখা", 4.8, 5.45, 4.1, 0.3, size=11, color=(255, 159, 28))

prs.save(OUT_PRENORM)
print(OUT_PRENORM)
