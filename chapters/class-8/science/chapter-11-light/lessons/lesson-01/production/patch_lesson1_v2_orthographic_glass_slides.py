"""Patch Lesson 1 v2 slides 26-27 with an orthographic realistic glass slab.

Purpose:
- Fix the user-identified problem that the previous 3D perspective glass slab
  looked realistic but the 2D ray overlay did not share the same optical plane.
- Keep the physical glass visually realistic while restoring textbook-correct
  ray-geometry clarity.
- Preserve Bangla-first textbook terminology.

Input:
  Class8_Science_Ch11_Lesson1_v2_Bangla_Textbook_RealGlass_PPTX_SAFE.pptx

Output:
  Class8_Science_Ch11_Lesson1_v2_Bangla_Textbook_OrthoGlass_PPTX_SAFE.pptx

Implementation note:
The slab asset is a straight-on/orthographic photorealistic glass image.
The ray path, interfaces, normals and labels are deterministic overlays.
"""

from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN
import hashlib
import math

BASE = Path('/mnt/data/science_slides_l1_v2/Class8_Science_Ch11_Lesson1_v2_Bangla_Textbook_RealGlass_PPTX_SAFE.pptx')
IMG = Path('/mnt/data/a_clean_minimal_still_life_industrial_product_sty.png')
WORK = Path('/mnt/data/science_slides_l1_v2/orthoglass')
WORK.mkdir(exist_ok=True)
CROP = WORK / 'orthographic_glass_slab_crop.png'
OUT_PRE = Path('/mnt/data/science_slides_l1_v2/Class8_Science_Ch11_Lesson1_v2_Bangla_Textbook_OrthoGlass_PRENORM.pptx')

im = Image.open(IMG).convert('RGBA')
w, h = im.size
im.crop((150, 245, w - 145, 770)).save(CROP)

prs = Presentation(BASE)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BG = RGBColor(7, 20, 38)
CYAN = RGBColor(155, 234, 254)
WHITE = RGBColor(246, 251, 255)
SOFT = RGBColor(183, 205, 232)
MUTED = RGBColor(127, 167, 204)
YELLOW = RGBColor(255, 209, 102)
GREEN = RGBColor(46, 229, 157)
PANEL = RGBColor(16, 44, 77)
GREEN_DARK = RGBColor(7, 65, 49)
ORANGE_DARK = RGBColor(55, 34, 10)
FONT = 'Noto Sans Bengali'


def clear_slide(slide):
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)


def fill_bg(slide):
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = BG
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08))
    top.fill.solid(); top.fill.fore_color.rgb = RGBColor(26, 126, 151); top.line.fill.background()
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.42), SLIDE_W, Inches(0.08))
    bottom.fill.solid(); bottom.fill.fore_color.rgb = RGBColor(26, 126, 151); bottom.line.fill.background()
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.55), Inches(6.86), Inches(12.8), Inches(6.86))
    line.line.color.rgb = RGBColor(32, 96, 135); line.line.width = Pt(1)


def add_text(slide, text, x, y, w, h, size=16, color=WHITE, bold=False, align='left', fill=None, line=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        box.fill.solid(); box.fill.fore_color.rgb = fill
    if line:
        box.line.color.rgb = line; box.line.width = Pt(1.2)
    else:
        box.line.fill.background()
    tf = box.text_frame; tf.clear()
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    run = p.add_run(); run.text = text
    run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return box


def add_badge(slide, text, x, y, w, h=0.45, color=CYAN, fill_rgb=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill_rgb or RGBColor(15, 42, 72)
    shp.line.color.rgb = color; shp.line.width = Pt(1.2)
    tf = shp.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.name = FONT; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = color
    return shp


def add_line(slide, x1, y1, x2, y2, color, width=1.4, dash=False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color; line.line.width = Pt(width)
    if dash:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def add_arrow(slide, x1, y1, x2, y2, color, width=4, dash=False):
    line = add_line(slide, x1, y1, x2, y2, color, width, dash)
    dx, dy = x2 - x1, y2 - y1
    angle = math.degrees(math.atan2(dy, dx))
    tri_w, tri_h = 0.20, 0.18
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x2 - tri_w / 2), Inches(y2 - tri_h / 2), Inches(tri_w), Inches(tri_h))
    tri.fill.solid(); tri.fill.fore_color.rgb = color
    tri.line.color.rgb = color
    tri.rotation = angle + 90
    return line


def add_label(slide, text, x, y, w, h, color=WHITE, fill=PANEL, line_color=CYAN, size=12, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color; shp.line.width = Pt(1.1)
    tf = shp.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return shp


def add_glass_panel(slide, x=2.55, y=2.05, w=8.2, h=3.15):
    slide.shapes.add_picture(str(CROP), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    left_x, right_x = 3.45, 9.85
    add_line(slide, left_x, 2.32, left_x, 4.86, CYAN, 1.5, False)
    add_line(slide, right_x, 2.32, right_x, 4.86, CYAN, 1.5, False)
    return left_x, right_x


def add_slide26(slide):
    clear_slide(slide); fill_bg(slide)
    add_text(slide, 'ঐচ্ছিক • কাচের পাত', 0.55, 0.2, 3.0, 0.25, 9, CYAN, True)
    add_text(slide, 'L1-F01a', 11.85, 0.2, 0.9, 0.25, 8, MUTED, False, 'right')
    add_text(slide, 'সমান্তরাল-পৃষ্ঠ কাচের পাত — ধাপ ১: পথ অনুমান', 0.8, 0.82, 10.7, 0.62, 26, WHITE, True)
    add_badge(slide, 'ধাপ ১ — আগে অনুমান', 10.2, 0.82, 2.2, 0.42, YELLOW, ORANGE_DARK)
    add_text(slide, 'বিভেদতলে আপতিত রশ্মি কোন দিকে বেঁকবে? দ্বিতীয় বিভেদতল দিয়ে বের হলে কী হবে?', 0.83, 1.52, 11.8, 0.38, 14, SOFT, False)
    left_x, right_x = add_glass_panel(slide)
    entry = (left_x, 3.68)
    add_arrow(slide, 0.95, 2.78, entry[0], entry[1], YELLOW, 4.2)
    add_line(slide, 2.55, entry[1], 5.35, entry[1], SOFT, 1.2, True)
    add_label(slide, 'আপতিত রশ্মি', 0.7, 2.43, 1.35, 0.34, YELLOW, ORANGE_DARK, YELLOW, 11.5)
    add_label(slide, 'প্রথম বিভেদতল', left_x - 0.9, 5.03, 1.8, 0.34, SOFT, RGBColor(15, 42, 72), CYAN, 10.5, False)
    add_label(slide, 'দ্বিতীয় বিভেদতল', right_x - 0.95, 5.03, 1.9, 0.34, SOFT, RGBColor(15, 42, 72), CYAN, 10.5, False)
    add_label(slide, 'বাস্তব কাচের পাত', 5.45, 2.08, 1.65, 0.34, CYAN, RGBColor(15, 42, 72), CYAN, 11.5)
    add_label(slide, 'অভিলম্ব', 4.05, 3.33, 0.95, 0.3, SOFT, RGBColor(15, 42, 72), SOFT, 10.5, False)
    add_text(slide, 'প্রকাশের আগে অনুমান: নির্গত রশ্মি কি আপতিত রশ্মির সাথে সমান্তরাল হবে?', 1.8, 6.02, 9.8, 0.5, 15, YELLOW, True, 'center', ORANGE_DARK, YELLOW)
    add_text(slide, 'ক্লিক/পরের ধাপ', 11.0, 6.94, 1.4, 0.25, 8.5, MUTED, False, 'right')


def add_slide27(slide):
    clear_slide(slide); fill_bg(slide)
    add_text(slide, 'ঐচ্ছিক • কাচের পাত', 0.55, 0.2, 3.0, 0.25, 9, CYAN, True)
    add_text(slide, 'L1-F01b', 11.85, 0.2, 0.9, 0.25, 8, MUTED, False, 'right')
    add_text(slide, 'সমান্তরাল-পৃষ্ঠ কাচের পাত — ধাপ ২: উত্তর ও ব্যাখ্যা', 0.8, 0.82, 10.6, 0.62, 25, WHITE, True)
    add_badge(slide, 'ধাপ ২ — উত্তর', 10.45, 0.84, 1.9, 0.42, GREEN, GREEN_DARK)
    add_text(slide, 'এই সমতা সমান্তরাল-পৃষ্ঠ কাচের পাতের জ্যামিতির জন্য; সব বিভেদতলে নয়।', 0.83, 1.52, 11.8, 0.36, 13.5, SOFT, False)
    left_x, right_x = add_glass_panel(slide)
    entry = (left_x, 3.68); exit = (right_x, 4.26)
    add_line(slide, 2.55, entry[1], 5.35, entry[1], SOFT, 1.1, True)
    add_line(slide, 8.15, exit[1], 11.0, exit[1], SOFT, 1.1, True)
    add_arrow(slide, 0.95, 2.78, entry[0], entry[1], YELLOW, 4.2)
    add_arrow(slide, entry[0], entry[1], exit[0], exit[1], GREEN, 4.2)
    add_arrow(slide, exit[0], exit[1], 12.35, 5.10, YELLOW, 4.2)
    add_label(slide, 'আপতিত রশ্মি', 0.66, 2.42, 1.35, 0.34, YELLOW, ORANGE_DARK, YELLOW, 11.5)
    add_label(slide, 'প্রতিসরিত রশ্মি', 5.25, 3.33, 1.45, 0.34, GREEN, RGBColor(8, 62, 47), GREEN, 11.5)
    add_label(slide, 'নির্গত রশ্মি', 10.55, 4.88, 1.3, 0.34, YELLOW, ORANGE_DARK, YELLOW, 11.5)
    add_label(slide, 'প্রথম বিভেদতল', left_x - 0.92, 5.03, 1.84, 0.34, SOFT, RGBColor(15, 42, 72), CYAN, 10.5, False)
    add_label(slide, 'দ্বিতীয় বিভেদতল', right_x - 0.95, 5.03, 1.9, 0.34, SOFT, RGBColor(15, 42, 72), CYAN, 10.5, False)
    add_label(slide, 'বাস্তব কাচের পাত', 5.45, 2.08, 1.65, 0.34, CYAN, RGBColor(15, 42, 72), CYAN, 11.5)
    add_label(slide, 'অভিলম্ব', 4.0, 3.32, 0.95, 0.3, SOFT, RGBColor(15, 42, 72), SOFT, 10.5, False)
    add_label(slide, 'অভিলম্ব', 8.55, 4.52, 0.95, 0.3, SOFT, RGBColor(15, 42, 72), SOFT, 10.5, False)
    add_text(slide, 'জ্যামিতিনির্ভর নিয়ম: আপতিত রশ্মি ও নির্গত রশ্মি সমান্তরাল; পার্শ্ব সরণ থাকতে পারে।', 1.25, 6.02, 10.75, 0.5, 14.2, GREEN, True, 'center', GREEN_DARK, GREEN)


add_slide26(prs.slides[25])
add_slide27(prs.slides[26])
prs.save(OUT_PRE)
print(OUT_PRE)
print('crop', CROP)
print('sha pre', hashlib.sha256(OUT_PRE.read_bytes()).hexdigest())
